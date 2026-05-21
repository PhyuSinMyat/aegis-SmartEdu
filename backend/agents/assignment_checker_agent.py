from __future__ import annotations

import io
import json as _json
import re
import zipfile
from pathlib import Path
from typing import Optional
from xml.etree import ElementTree as ET

from config import Config
from llm_service import generate_response


_PROMPT_PATH = Config.PROMPTS_DIR / "assignment_checker_prompt.txt"
_BEDROCK_MAX_NATIVE_BYTES = 4_500_000
_MAX_SUBMISSION_CHARS_FOR_MODEL = 180_000
_SUBMISSION_SAMPLE_CHARS = 12_000

# Images sent as vision blocks
_IMAGE_FORMAT_MAP = {
    ".png": "png",
    ".jpg": "jpeg",
    ".jpeg": "jpeg",
    ".webp": "webp",
    ".gif": "gif",
}

# Bedrock Converse native document formats — sent as raw bytes, Claude reads natively
_BEDROCK_NATIVE_FORMATS = {
    ".pdf":  "pdf",
    ".docx": "docx",
    ".doc":  "docx",
    ".xlsx": "xlsx",
    ".xls":  "xlsx",
    ".csv":  "csv",
    ".txt":  "txt",
    ".html": "html",
    ".htm":  "html",
    ".md":   "md",
}

# Code / data / markup files — read as UTF-8 text and sent inline
_CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".cpp", ".c", ".h", ".cs",
    ".php", ".rb", ".go", ".rs", ".swift", ".kt",
    ".r", ".m", ".sh", ".bash", ".ps1",
    ".sql", ".css", ".scss", ".sass", ".less",
    ".vue", ".dart", ".json", ".xml",
    ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".rst", ".tex", ".lua", ".scala", ".jl",
    ".vb", ".pl", ".elm", ".ex", ".exs",
    ".erl", ".zig", ".nim",
}

# Human-readable language labels for code blocks
_EXT_LANG: dict[str, str] = {
    ".py": "Python", ".js": "JavaScript", ".ts": "TypeScript",
    ".jsx": "JavaScript (JSX)", ".tsx": "TypeScript (TSX)",
    ".java": "Java", ".cpp": "C++", ".c": "C", ".h": "C/C++ Header",
    ".cs": "C#", ".php": "PHP", ".rb": "Ruby", ".go": "Go", ".rs": "Rust",
    ".swift": "Swift", ".kt": "Kotlin", ".r": "R",
    ".m": "MATLAB / Objective-C", ".sh": "Shell Script",
    ".bash": "Bash", ".ps1": "PowerShell",
    ".sql": "SQL", ".html": "HTML", ".htm": "HTML", ".css": "CSS",
    ".scss": "SCSS", ".sass": "Sass", ".less": "Less",
    ".vue": "Vue", ".dart": "Dart",
    ".json": "JSON", ".xml": "XML",
    ".yaml": "YAML", ".yml": "YAML",
    ".toml": "TOML", ".ini": "INI Config", ".cfg": "Config",
    ".md": "Markdown", ".rst": "reStructuredText", ".tex": "LaTeX",
    ".lua": "Lua", ".scala": "Scala", ".jl": "Julia",
    ".vb": "Visual Basic", ".pl": "Perl",
    ".elm": "Elm", ".ex": "Elixir", ".exs": "Elixir Script",
    ".erl": "Erlang", ".zig": "Zig", ".nim": "Nim",
}


def _load_system_prompt() -> str:
    if not _PROMPT_PATH.exists():
        raise FileNotFoundError(f"Assignment checker prompt not found: {_PROMPT_PATH}")
    return _PROMPT_PATH.read_text(encoding="utf-8").strip()


def _extract_json(raw: str) -> dict:
    raw = raw.strip()
    fence_match = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", raw)
    if fence_match:
        raw = fence_match.group(1).strip()
    start = raw.find("{")
    end = raw.rfind("}")
    if start == -1 or end == -1:
        raise ValueError("No JSON object found in LLM response.")
    return _json.loads(raw[start : end + 1])


def _make_image_block(image_bytes: bytes, fmt: str) -> dict:
    return {"image": {"format": fmt, "source": {"bytes": image_bytes}}}


def _make_doc_block(file_bytes: bytes, fmt: str, filename: str = "") -> dict:
    stem = Path(filename).stem if filename else "document"
    # Bedrock doc name: alphanumeric, spaces, hyphens, underscores only
    name = re.sub(r"[^a-zA-Z0-9 _\-]", "_", stem)[:60].strip("_") or "document"
    return {
        "document": {
            "name": name,
            "format": fmt,
            "source": {"bytes": file_bytes},
        }
    }


class AssignmentCheckerAgent:
    """
    Analyses a student assignment submission and returns structured feedback.

    Accepts virtually any file type:
      - Office / document: PDF, DOCX, XLSX (sent natively to Bedrock)
      - Notebook: .ipynb  (cells parsed to text)
      - Presentation: .pptx (slides parsed to text)
      - Images: PNG, JPG, WEBP, GIF (vision blocks)
      - Code / data: .py, .js, .java, .sql, .json, .yaml, etc. (UTF-8 text)
    """

    STRICTNESS_LEVELS = {"lenient", "normal", "strict"}

    @classmethod
    def prepare_submission_for_model(cls, assignment_text: str) -> str:
        text = assignment_text.strip()
        if len(text) <= _MAX_SUBMISSION_CHARS_FOR_MODEL:
            return text

        sample = _SUBMISSION_SAMPLE_CHARS
        chunks: list[str] = []
        chunks.append(
            "[Large submission note]\n"
            f"The extracted submission contains {len(text):,} characters. "
            "The app accepts the full upload, but this AI request includes representative excerpts "
            "from the beginning, middle, and end so the model stays within its context limit. "
            "For precise grading of very large projects, review file-specific feedback against the original files.\n"
        )
        chunks.append("[Beginning excerpt]\n" + text[:sample])

        middle_start = max(0, (len(text) // 2) - (sample // 2))
        chunks.append("[Middle excerpt]\n" + text[middle_start : middle_start + sample])

        chunks.append("[Ending excerpt]\n" + text[-sample:])

        file_headers = re.findall(r"^SUBMISSION FILE \d+:[^\n]+", text, flags=re.MULTILINE)
        if file_headers:
            chunks.append("[Uploaded file index]\n" + "\n".join(file_headers[:200]))

        return "\n\n---\n\n".join(chunks)

    @classmethod
    def analyse(
        cls,
        assignment_text: str = "",
        strictness: str = "strict",
        brief_text: Optional[str] = None,
        requirements: Optional[str] = None,
        instructions_text: Optional[str] = None,
        submission_image_blocks: Optional[list] = None,
        brief_image_blocks: Optional[list] = None,
        instructions_image_blocks: Optional[list] = None,
    ) -> dict:
        strictness = strictness.lower().strip()
        if strictness not in cls.STRICTNESS_LEVELS:
            strictness = "strict"

        parts = [f"STRICTNESS LEVEL: {strictness.upper()}"]

        # ── Project / Assignment Instructions (optional) ──────────────────────
        if instructions_text and instructions_text.strip():
            instr = "PROJECT/ASSIGNMENT INSTRUCTIONS (what the student was asked to do):\n" + instructions_text.strip()
            if instructions_image_blocks:
                instr += "\n\n[The instructions document is also attached above.]"
            parts.append(instr)
        elif instructions_image_blocks:
            parts.append(
                "PROJECT/ASSIGNMENT INSTRUCTIONS: "
                "[Attached document/image above — extract the task description and requirements from it.]"
            )

        # ── Rubric / Brief (required) ─────────────────────────────────────────
        if brief_text and brief_text.strip():
            rubric = "ASSIGNMENT RUBRIC/BRIEF (grading criteria — REQUIRED):\n" + brief_text.strip()
            if brief_image_blocks:
                rubric += "\n\n[The rubric document is also attached above.]"
            parts.append(rubric)
        elif brief_image_blocks:
            parts.append(
                "ASSIGNMENT RUBRIC/BRIEF (grading criteria — REQUIRED): "
                "[Attached document/image above — extract all rubric criteria, marks, and grade bands from it.]"
            )
        else:
            parts.append(
                "ASSIGNMENT RUBRIC/BRIEF: [NOT PROVIDED] — "
                "Return a JSON error: {\"error\": \"Rubric not provided. Upload the rubric before grading.\"}"
            )

        # ── Student Submission ────────────────────────────────────────────────
        if assignment_text.strip():
            sub = "STUDENT SUBMISSION:\n" + cls.prepare_submission_for_model(assignment_text)
            if submission_image_blocks:
                sub += "\n\n[The submission document is also attached above.]"
            parts.append(sub)
        elif submission_image_blocks:
            parts.append(
                "STUDENT SUBMISSION: "
                "[Attached document/image above — evaluate the content visible in it.]"
            )
        else:
            parts.append("STUDENT SUBMISSION: [No content provided]")

        user_message = "\n\n---\n\n".join(parts)
        system_prompt = _load_system_prompt()

        # Build content blocks — doc/image blocks first, text message last
        content_blocks: list = []
        if instructions_image_blocks:
            content_blocks.extend(instructions_image_blocks)
        if brief_image_blocks:
            content_blocks.extend(brief_image_blocks)
        if submission_image_blocks:
            content_blocks.extend(submission_image_blocks)
        content_blocks.append({"text": user_message})

        raw_response = generate_response(
            content_blocks=content_blocks,
            system_prompt=system_prompt,
            max_tokens=8000,
        )
        return _extract_json(raw_response)

    # ── Content extraction ─────────────────────────────────────────────────────

    @classmethod
    def extract_content_from_file(
        cls, file_bytes: bytes, suffix: str, filename: str = ""
    ) -> tuple[str, list]:
        """
        Dispatch to the correct extractor.
        Returns (text, content_blocks) where content_blocks are Bedrock content dicts
        (image blocks or native document blocks).
        """
        # Native Bedrock document formats — best quality, Claude reads the file directly
        if suffix in _BEDROCK_NATIVE_FORMATS:
            if len(file_bytes) > _BEDROCK_MAX_NATIVE_BYTES:
                return cls.extract_large_native_file_as_text(file_bytes, suffix)
            fmt = _BEDROCK_NATIVE_FORMATS[suffix]
            return "", [_make_doc_block(file_bytes, fmt, filename)]

        # Jupyter notebook — parse cells
        if suffix == ".ipynb":
            return cls.extract_text_from_ipynb(file_bytes)

        # PowerPoint — parse slides
        if suffix in (".pptx", ".ppt"):
            return cls.extract_text_from_pptx(file_bytes)

        # Images — vision blocks
        if suffix in _IMAGE_FORMAT_MAP:
            return cls.extract_image_file(file_bytes, suffix)

        # Code / data / text files — UTF-8 text with language label
        if suffix in _CODE_EXTENSIONS:
            return cls.extract_code_file(file_bytes, suffix)

        raise ValueError(f"Unsupported file type: '{suffix}'")

    @classmethod
    def extract_large_native_file_as_text(cls, file_bytes: bytes, suffix: str) -> tuple[str, list]:
        if suffix == ".pdf":
            return cls.extract_pdf_file_as_text(file_bytes), []
        if suffix in (".xlsx", ".xls"):
            return cls.extract_spreadsheet_file_as_text(file_bytes), []
        if suffix == ".docx":
            return cls.extract_docx_file_as_text(file_bytes), []
        if suffix == ".doc":
            raise ValueError(
                "DOC files larger than 4.5 MB cannot be read directly. "
                "Please convert the file to DOCX or PDF and try again."
            )
        if suffix in {".csv", ".txt", ".html", ".htm", ".md"}:
            return cls.extract_code_file(file_bytes, suffix)
        raise ValueError(f"Unsupported large document type: '{suffix}'")

    @classmethod
    def extract_pdf_file_as_text(cls, file_bytes: bytes) -> str:
        try:
            import pdfplumber
        except ImportError as exc:
            raise ValueError("Install pdfplumber to read PDF files larger than 4.5 MB.") from exc

        parts: list[str] = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(f"[Page {page_num}]\n{text.strip()}")
        return "\n\n".join(parts).strip()

    @classmethod
    def extract_spreadsheet_file_as_text(cls, file_bytes: bytes) -> str:
        try:
            import openpyxl
        except ImportError as exc:
            raise ValueError("Install openpyxl to read Excel files larger than 4.5 MB.") from exc

        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
        parts: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    rows.append("\t".join(values).rstrip())
            if rows:
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        workbook.close()
        return "\n\n".join(parts).strip()

    @classmethod
    def extract_docx_file_as_text(cls, file_bytes: bytes) -> str:
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as docx:
                xml_bytes = docx.read("word/document.xml")
        except Exception as exc:
            raise ValueError("Could not read DOCX text. Try exporting the document as PDF.") from exc

        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        root = ET.fromstring(xml_bytes)
        paragraphs: list[str] = []
        for para in root.findall(".//w:p", ns):
            chunks = [node.text or "" for node in para.findall(".//w:t", ns)]
            text = "".join(chunks).strip()
            if text:
                paragraphs.append(text)
        return "\n\n".join(paragraphs).strip()

    @classmethod
    def extract_text_from_ipynb(cls, file_bytes: bytes) -> tuple[str, list]:
        try:
            nb = _json.loads(file_bytes.decode("utf-8"))
        except Exception:
            return file_bytes.decode("utf-8", errors="replace"), []

        parts: list[str] = []
        for i, cell in enumerate(nb.get("cells", []), 1):
            cell_type = cell.get("cell_type", "")
            source = "".join(cell.get("source", []))
            if not source.strip():
                continue
            if cell_type == "markdown":
                parts.append(f"[Markdown Cell {i}]\n{source}")
            elif cell_type == "code":
                outputs: list[str] = []
                for out in cell.get("outputs", []):
                    if out.get("output_type") in ("stream", "execute_result", "display_data"):
                        text = out.get("text") or out.get("data", {}).get("text/plain", [])
                        if isinstance(text, list):
                            text = "".join(text)
                        if isinstance(text, str) and text.strip():
                            outputs.append(text.strip())
                cell_text = f"[Code Cell {i}]\n{source}"
                if outputs:
                    cell_text += "\n[Output]\n" + "\n".join(outputs)
                parts.append(cell_text)
            else:
                parts.append(f"[Cell {i} ({cell_type})]\n{source}")

        return "\n\n".join(parts).strip(), []

    @classmethod
    def extract_text_from_pptx(cls, file_bytes: bytes) -> tuple[str, list]:
        try:
            from pptx import Presentation
        except ImportError:
            return "[PowerPoint file — install python-pptx to extract content]", []

        prs = Presentation(io.BytesIO(file_bytes))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
            if lines:
                parts.append(f"[Slide {i}]\n" + "\n".join(lines))
        return "\n\n".join(parts).strip(), []

    @classmethod
    def extract_image_file(cls, file_bytes: bytes, suffix: str) -> tuple[str, list]:
        if len(file_bytes) > _BEDROCK_MAX_NATIVE_BYTES:
            file_bytes, suffix = cls.compress_image_for_bedrock(file_bytes)
        fmt = _IMAGE_FORMAT_MAP.get(suffix, "jpeg")
        return "", [_make_image_block(file_bytes, fmt)]

    @classmethod
    def compress_image_for_bedrock(cls, file_bytes: bytes) -> tuple[bytes, str]:
        try:
            from PIL import Image
        except ImportError as exc:
            raise ValueError("Install Pillow to process image files larger than 4.5 MB.") from exc

        with Image.open(io.BytesIO(file_bytes)) as img:
            if getattr(img, "is_animated", False):
                img.seek(0)
            image = img.convert("RGB")

        for quality in (85, 75, 65, 55, 45, 35):
            out = io.BytesIO()
            image.save(out, format="JPEG", quality=quality, optimize=True)
            compressed = out.getvalue()
            if len(compressed) <= _BEDROCK_MAX_NATIVE_BYTES:
                return compressed, ".jpg"

        scale = 0.85
        while scale >= 0.25:
            resized = image.resize(
                (max(1, int(image.width * scale)), max(1, int(image.height * scale)))
            )
            out = io.BytesIO()
            resized.save(out, format="JPEG", quality=70, optimize=True)
            compressed = out.getvalue()
            if len(compressed) <= _BEDROCK_MAX_NATIVE_BYTES:
                return compressed, ".jpg"
            scale -= 0.15

        raise ValueError("Image is too large to process. Please compress it and try again.")

    @classmethod
    def extract_code_file(cls, file_bytes: bytes, suffix: str) -> tuple[str, list]:
        try:
            text = file_bytes.decode("utf-8-sig")
        except UnicodeDecodeError:
            try:
                text = file_bytes.decode("latin-1")
            except Exception:
                return "[Binary file — cannot display as text]", []
        lang = _EXT_LANG.get(suffix, "Text")
        return f"[{lang} File]\n{text.strip()}", []

    # ── Legacy helpers kept for compatibility ──────────────────────────────────

    @classmethod
    def extract_text_from_file(cls, file_bytes: bytes, suffix: str) -> str:
        text, _ = cls.extract_content_from_file(file_bytes, suffix)
        return text

    @classmethod
    def extract_text_from_pdf(cls, file_bytes: bytes) -> tuple[str, list]:
        return cls.extract_content_from_file(file_bytes, ".pdf")

    @classmethod
    def extract_text_from_docx(cls, file_bytes: bytes) -> tuple[str, list]:
        return cls.extract_content_from_file(file_bytes, ".docx")

    @classmethod
    def extract_text_from_xlsx(cls, file_bytes: bytes) -> tuple[str, list]:
        return cls.extract_content_from_file(file_bytes, ".xlsx")

    @classmethod
    def extract_text_from_txt(cls, file_bytes: bytes) -> tuple[str, list]:
        return cls.extract_content_from_file(file_bytes, ".txt")
